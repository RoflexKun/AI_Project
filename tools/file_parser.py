
import os
import fitz

from pathlib import Path
from pptx import Presentation

class FileParser:
    output_path = ''
    file_folder_dict = {}

    def __init__(self, file_path):
        self.file_path = Path(file_path)

    def create_output_directory(self, script_path):
        parent_path = script_path[ : script_path.rfind('\\')]

        FileParser.output_path = parent_path + '\\output'
        os.makedirs(FileParser.output_path, exist_ok=True)

    def parse_files(self):
        for file in self.file_path.iterdir():
            file_name = file.name
            if file_name.endswith('.pptx'):
                self.parse_pptx(file)
            elif file_name.endswith('.pdf'):
                self.parse_pdf(file)

    def parse_pptx(self, file_path):
        powerpoint_file = Presentation(file_path)
        powerpoint_name = file_path.name[ : file_path.name.rfind('.')]

        powerpoint_folder = FileParser.output_path + '\\' + powerpoint_name
        os.makedirs(powerpoint_folder, exist_ok=True)
        FileParser.file_folder_dict[powerpoint_name] = powerpoint_folder
        slides_folder = powerpoint_folder + '\\slides'
        os.makedirs(slides_folder, exist_ok=True)

        slide_index = 1
        for slide in powerpoint_file.slides:
            slide_folder = Path(slides_folder) / f"slide{slide_index}"
            os.makedirs(slide_folder, exist_ok=True)

            output_file_path = slide_folder / f"slide{slide_index}_text.txt"

            photos_folder = None

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    if not Path(output_file_path).exists():
                        Path(output_file_path).write_text(shape.text_frame.text, encoding="utf-8")
                    else:
                        with open(output_file_path, 'a', encoding='utf-8') as output_file:
                            output_file.write(shape.text_frame.text)

                if shape.shape_type == 13:
                    if photos_folder is None:
                        photos_folder = slide_folder / "photos"
                        os.makedirs(photos_folder, exist_ok=True)

                    image = shape.image
                    image_extension = image.ext
                    image_name = photos_folder / f"image_{slide_index}.{image_extension}"
                    with open(image_name, 'wb') as image_file:
                        image_file.write(image.blob)

            slide_index += 1

    def parse_pdf(self, file_path):
        pdf = fitz.open(file_path)

        pdf_name = file_path.name[ : file_path.name.rfind('.')]

        pdf_folder = FileParser.output_path + '\\' + pdf_name
        os.makedirs(pdf_folder, exist_ok=True)
        FileParser.file_folder_dict[pdf_name] = pdf_folder

        pages_folder = pdf_folder + '\\pages'
        os.makedirs(pages_folder, exist_ok=True)

        page_index = 1
        for page in pdf:
            page_folder = Path(pages_folder) / f"page{page_index}"
            os.makedirs(page_folder, exist_ok=True)

            output_file_path = page_folder / f"page{page_index}_text.txt"

            text = page.get_text()
            Path(output_file_path).write_text(text, encoding="utf-8")

            photo_folder = None

            images = page.get_images(full=True)
            for image in images:
                xref = image[0]
                pix = fitz.Pixmap(pdf, xref)

                if photo_folder is None:
                    photo_folder = page_folder / "photos"
                    os.makedirs(photo_folder, exist_ok=True)

                image_extension = 'png' if pix.alpha or pix.n > 3 else 'jpg'
                image_name = photo_folder / f"image_{page_index}.{image_extension}"

                pix.save(str(image_name))

            page_index += 1
